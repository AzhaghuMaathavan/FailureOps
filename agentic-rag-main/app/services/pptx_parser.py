import uuid
import logging
from app.models.document import Page, DocumentBlock

logger = logging.getLogger(__name__)

def parse_pptx_to_blocks(file_path: str, doc_id: str, db):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except ImportError:
        logger.warning("python-pptx not installed. Skipping pptx parsing.")
        return False
    except Exception as e:
        logger.error(f"Failed to load PPTX {file_path}: {e}")
        return False

    for slide_idx, slide in enumerate(prs.slides):
        page_id = str(uuid.uuid4())
        db_page = Page(
            id=page_id,
            document_id=doc_id,
            page_number=slide_idx,
            image_path="",
            status="COMPLETED"
        )
        db.add(db_page)
        
        block_idx = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    db_block = DocumentBlock(
                        id=str(uuid.uuid4()),
                        document_id=doc_id,
                        page_id=page_id,
                        block_index=block_idx,
                        block_type="Paragraph",
                        content=text,
                        raw_metadata={"slide": [slide_idx + 1]}
                    )
                    db.add(db_block)
                    block_idx += 1
                    
            if shape.has_table:
                table = shape.table
                lines = []
                for r_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    lines.append(f"| {' | '.join(cells)} |")
                    if r_idx == 0:
                        lines.append(f"| {' | '.join(['---'] * len(cells))} |")
                        
                table_md = "\n".join(lines)
                db_block = DocumentBlock(
                    id=str(uuid.uuid4()),
                    document_id=doc_id,
                    page_id=page_id,
                    block_index=block_idx,
                    block_type="Table",
                    content=table_md,
                    raw_metadata={"slide": [slide_idx + 1]}
                )
                db.add(db_block)
                block_idx += 1

    db.commit()
    return True
